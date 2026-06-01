"""ApplyOS Investor-Pitchdeck Generator.

Erzeugt dist/ApplyOS-Pitchdeck.pptx im 16:9-Format mit dem Variant-E-Twilight
Design-System der Marketing-Website. 12 Slides.

Run: python3 scripts/build_pitchdeck.py
"""

from __future__ import annotations

import os
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ──────────────────────────────────────────────────────────────────────────────
# Paths
# ──────────────────────────────────────────────────────────────────────────────

REPO_ROOT = Path(__file__).resolve().parent.parent
ASSET_LOGO = REPO_ROOT / "public" / "appi-logo.png"
ASSET_RESULT = REPO_ROOT / "public" / "app-result.png"
ASSET_ICON = REPO_ROOT / "public" / "icon-512.png"
OUTPUT_PATH = REPO_ROOT / "dist" / "ApplyOS-Pitchdeck.pptx"

# ──────────────────────────────────────────────────────────────────────────────
# Design-Tokens (Variant E "Twilight")
# ──────────────────────────────────────────────────────────────────────────────

# Colors
BG = RGBColor(0x1C, 0x1A, 0x17)        # Body background (warm graphite)
PAPER = RGBColor(0x2A, 0x26, 0x20)      # Card background
PAPER_HI = RGBColor(0x3A, 0x34, 0x2C)   # Card hover/higher contrast
DEEP = RGBColor(0x15, 0x13, 0x0F)       # Footer/deep zones
TEXT = RGBColor(0xED, 0xE4, 0xD0)       # Primary text
TEXT2 = RGBColor(0xD4, 0xC9, 0xB0)      # Secondary text
MUTED = RGBColor(0xA3, 0x97, 0x82)      # Muted text
FAINT = RGBColor(0x6B, 0x62, 0x53)      # Faint labels
TEAL = RGBColor(0x3A, 0xAB, 0x83)       # Primary accent (Jobsuchende)
TAN = RGBColor(0xC9, 0xA5, 0x87)        # Secondary accent (Schüler)
ERROR = RGBColor(0xD7, 0x7A, 0x7A)      # Problem indicators
HAIRLINE_BG = RGBColor(0x2D, 0x2A, 0x26)  # Approximated 8% white on BG
HAIRLINE_CARD = RGBColor(0x3A, 0x35, 0x2E)  # Approximated 8% white on Card
TEAL_GLOW = RGBColor(0x23, 0x32, 0x2A)  # Subtle teal tint (10% teal on BG)
TAN_GLOW = RGBColor(0x32, 0x2D, 0x25)   # Subtle tan tint

# Fonts
FONT_SERIF = "Source Serif 4"      # Fallback: Georgia
FONT_SERIF_ALT = "Georgia"
FONT_SANS = "Geist"                # Fallback: Helvetica Neue
FONT_SANS_ALT = "Helvetica Neue"

# Slide dimensions (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.6)
MARGIN_TOP = Inches(0.55)
MARGIN_BOTTOM = Inches(0.4)
TOTAL_SLIDES = 12

# ──────────────────────────────────────────────────────────────────────────────
# Low-level XML Helpers
# ──────────────────────────────────────────────────────────────────────────────


def set_letter_spacing(run, spc_hundredths_pt: int) -> None:
    """Setzt letter-spacing (spc) am Run via XML. Wert in 1/100 pt."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(spc_hundredths_pt))


def set_fill_alpha(shape, alpha_pct: int) -> None:
    """Setzt Fill-Alpha (0-100) via XML auf das solidFill-Element."""
    spPr = shape.fill._xPr.find(qn("a:solidFill"))
    if spPr is None:
        return
    srgb = spPr.find(qn("a:srgbClr"))
    if srgb is None:
        return
    # Entferne bestehende alpha-Elemente
    for existing in srgb.findall(qn("a:alpha")):
        srgb.remove(existing)
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", str(int(alpha_pct * 1000)))


def remove_shadow(shape) -> None:
    """Entfernt Standard-Schatten von der Shape."""
    spPr = shape._element.spPr
    # Add empty effectLst that overrides theme-level shadow
    for existing in spPr.findall(qn("a:effectLst")):
        spPr.remove(existing)
    etree.SubElement(spPr, qn("a:effectLst"))


def set_text_anchor(text_frame, vertical: str = "top") -> None:
    """Setzt vertical anchor (top, middle, bottom)."""
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    text_frame.vertical_anchor = anchor_map[vertical]


# ──────────────────────────────────────────────────────────────────────────────
# Mid-level Helpers (Background, Shapes, Text)
# ──────────────────────────────────────────────────────────────────────────────


def new_slide(prs):
    """Erstellt eine neue Blank-Slide mit Dark-Background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    return slide


def add_dark_background(slide) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    remove_shadow(bg)


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None,
             border_width_pt=0.5, radius_pt=None):
    """Fügt ein Rectangle hinzu (rounded wenn radius_pt gesetzt)."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_pt else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color is not None:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width_pt)
    else:
        shape.line.fill.background()
    remove_shadow(shape)
    # Set corner radius via adjustment value (relative to shorter side)
    if radius_pt and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        # Adjustment value (0.0-0.5) controls radius proportional to shorter dimension
        shorter = min(width, height)
        # radius_pt in points → convert to EMU
        radius_emu = Pt(radius_pt)
        adj_val = min(0.5, radius_emu / shorter)
        shape.adjustments[0] = adj_val
    return shape


def add_card(slide, left, top, width, height, fill_color=PAPER,
             border_color=HAIRLINE_CARD, radius_pt=12):
    """Standard-Card mit Paper-Fill und Hairline-Border."""
    return add_rect(slide, left, top, width, height, fill_color=fill_color,
                    border_color=border_color, border_width_pt=0.75, radius_pt=radius_pt)


def add_hairline(slide, left, top, width, color=HAIRLINE_BG, width_pt=0.5):
    """Dünne horizontale Linie."""
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    return line


def add_vertical_line(slide, left, top, height, color=HAIRLINE_BG, width_pt=0.5):
    line = slide.shapes.add_connector(1, left, top, left, top + height)
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    return line


def add_glow(slide, cx, cy, width, height, color=TEAL_GLOW):
    """Subtiler weicher Glow als Oval (approximiert ohne Soft-Edge)."""
    left = cx - width // 2
    top = cy - height // 2
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    glow.fill.solid()
    glow.fill.fore_color.rgb = color
    glow.line.fill.background()
    remove_shadow(glow)
    return glow


def add_text(slide, text, left, top, width, height, *,
             font=FONT_SANS, font_alt=FONT_SANS_ALT, size_pt=12,
             color=TEXT, bold=False, italic=False, align="left",
             anchor="top", letter_spacing=None, line_spacing=None):
    """Universeller Text-Helper."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    set_text_anchor(tf, anchor)
    p = tf.paragraphs[0]
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p.alignment = align_map[align]
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if letter_spacing is not None:
        set_letter_spacing(run, letter_spacing)
    # Set fallback font via XML (latin typeface) — pptx wählt automatisch wenn nicht installiert
    rPr = run._r.get_or_add_rPr()
    # Add East Asian/Complex script fallbacks (mostly cosmetic)
    return tb, run


def add_multiline_text(slide, lines, left, top, width, height, *,
                       font=FONT_SANS, size_pt=12, color=TEXT, bold=False,
                       italic=False, align="left", anchor="top",
                       line_spacing=1.4, paragraph_spacing_pt=4,
                       letter_spacing=None):
    """Mehrere Zeilen als separate Paragraphs."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    set_text_anchor(tf, anchor)
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(paragraph_spacing_pt)
        p.alignment = align_map[align]
        p.line_spacing = line_spacing
        if isinstance(line, dict):
            # Erweitertes Format: {'text': ..., 'italic': True/False, 'color': ..., 'size': ...}
            run = p.add_run()
            run.text = line["text"]
            run.font.name = line.get("font", font)
            run.font.size = Pt(line.get("size", size_pt))
            run.font.bold = line.get("bold", bold)
            run.font.italic = line.get("italic", italic)
            run.font.color.rgb = line.get("color", color)
            if letter_spacing is not None:
                set_letter_spacing(run, letter_spacing)
        else:
            run = p.add_run()
            run.text = line
            run.font.name = font
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            if letter_spacing is not None:
                set_letter_spacing(run, letter_spacing)
    return tb


# ──────────────────────────────────────────────────────────────────────────────
# High-level Composite Helpers
# ──────────────────────────────────────────────────────────────────────────────


def add_eyebrow(slide, text, left, top, color=TEAL, width=Inches(6)):
    """UPPERCASE Section-Label mit tracking."""
    tb, run = add_text(
        slide, text.upper(), left, top, width, Pt(20),
        font=FONT_SANS, size_pt=10, color=color, bold=True,
        letter_spacing=200,  # 0.14em bei ~14pt ≈ 200/100pt
    )
    return tb


def add_serif_heading(slide, text, left, top, width, height, *,
                      size_pt=44, color=TEXT, italic=False, align="left",
                      letter_spacing=-260):
    """Großer Serif-Heading (H1/H2)."""
    return add_text(
        slide, text, left, top, width, height,
        font=FONT_SERIF, size_pt=size_pt, color=color,
        italic=italic, align=align, letter_spacing=letter_spacing,
        line_spacing=1.05,
    )


def add_serif_heading_dual(slide, line1, line2, left, top, width, *,
                           size_pt=60, color1=TEXT, color2=TEAL,
                           letter_spacing=-360):
    """Zwei-Zeilen Serif-Heading (zweite Zeile italic + Accent)."""
    tb = slide.shapes.add_textbox(left, top, width, Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p1 = tf.paragraphs[0]
    p1.line_spacing = 1.0
    r1 = p1.add_run()
    r1.text = line1
    r1.font.name = FONT_SERIF
    r1.font.size = Pt(size_pt)
    r1.font.color.rgb = color1
    set_letter_spacing(r1, letter_spacing)

    p2 = tf.add_paragraph()
    p2.line_spacing = 1.0
    p2.space_before = Pt(2)
    r2 = p2.add_run()
    r2.text = line2
    r2.font.name = FONT_SERIF
    r2.font.size = Pt(size_pt)
    r2.font.italic = True
    r2.font.color.rgb = color2
    set_letter_spacing(r2, letter_spacing)
    return tb


def add_body(slide, text, left, top, width, height, *,
             size_pt=14, color=TEXT2, align="left", line_spacing=1.5):
    return add_text(
        slide, text, left, top, width, height,
        font=FONT_SANS, size_pt=size_pt, color=color, align=align,
        line_spacing=line_spacing,
    )


def add_stat(slide, number, label, left, top, *, width=Inches(3.0),
             number_color=TEAL, label_color=MUTED, number_size=54,
             align="left"):
    """Großer Stat-Block: italic-Teal-Zahl + kleines Label."""
    add_text(slide, number, left, top, width, Inches(1.0),
             font=FONT_SERIF, size_pt=number_size, color=number_color,
             italic=True, align=align, letter_spacing=-180)
    add_text(slide, label, left, top + Inches(0.95), width, Inches(0.4),
             font=FONT_SANS, size_pt=10, color=label_color,
             align=align, letter_spacing=180)


def add_footer(slide, page_num: int) -> None:
    """Footer mit Hairline + Page-Number + Brand."""
    footer_y = SLIDE_H - Inches(0.35)
    add_hairline(slide, MARGIN_X, footer_y - Inches(0.15),
                 SLIDE_W - 2 * MARGIN_X, color=HAIRLINE_BG)
    add_text(slide, f"ApplyOS · Vertraulich", MARGIN_X, footer_y - Inches(0.05),
             Inches(4), Inches(0.3),
             font=FONT_SANS, size_pt=8, color=FAINT,
             letter_spacing=120)
    add_text(slide, f"{page_num} / {TOTAL_SLIDES}",
             SLIDE_W - MARGIN_X - Inches(2), footer_y - Inches(0.05),
             Inches(2), Inches(0.3),
             font=FONT_SANS, size_pt=8, color=FAINT, align="right",
             letter_spacing=120)


def add_phone_frame(slide, left, top, width, height, *, screen_padding_pt=4):
    """Stilisierter iPhone-Frame mit Inner-Screen."""
    # Outer frame (dark bezel)
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0x15, 0x13, 0x0F)
    frame.line.color.rgb = HAIRLINE_CARD
    frame.line.width = Pt(0.75)
    frame.adjustments[0] = 0.10
    remove_shadow(frame)
    # Inner screen (App background = warm off-white #F4F3EF as on landing)
    pad = Pt(screen_padding_pt)
    screen_left = left + pad
    screen_top = top + pad
    screen_w = width - 2 * pad
    screen_h = height - 2 * pad
    screen = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, screen_left, screen_top, screen_w, screen_h
    )
    screen.fill.solid()
    screen.fill.fore_color.rgb = RGBColor(0xF4, 0xF3, 0xEF)
    screen.line.fill.background()
    screen.adjustments[0] = 0.085
    remove_shadow(screen)
    return frame, screen, (screen_left, screen_top, screen_w, screen_h)


def add_image_card(slide, image_path, left, top, width, height, radius_pt=14):
    """Bild mit Hairline-Border-Card drumherum."""
    # Card-Hintergrund
    bg = add_card(slide, left, top, width, height,
                  fill_color=PAPER, border_color=HAIRLINE_CARD, radius_pt=radius_pt)
    # Bild mittig einsetzen
    pad = Pt(8)
    img_left = left + pad
    img_top = top + pad
    img_w = width - 2 * pad
    img_h = height - 2 * pad
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), img_left, img_top,
                                 width=img_w, height=img_h)
    return bg


# ──────────────────────────────────────────────────────────────────────────────
# Slide Builders
# ──────────────────────────────────────────────────────────────────────────────


def build_slide_01_cover(prs):
    slide = new_slide(prs)
    # Hero-Glow zentriert
    add_glow(slide, SLIDE_W // 2, Inches(3.6),
             Inches(13), Inches(7), color=TEAL_GLOW)

    # Eyebrow (centered)
    add_text(slide, "INVESTOR DECK — 2026",
             Inches(0), Inches(1.3), SLIDE_W, Inches(0.4),
             font=FONT_SANS, size_pt=11, color=TEAL, bold=True,
             align="center", letter_spacing=280)

    # H1 dual-line (centered)
    add_text(slide, "Dein KI-Karriere-Agent.",
             Inches(0), Inches(2.1), SLIDE_W, Inches(1.4),
             font=FONT_SERIF, size_pt=66, color=TEXT,
             align="center", letter_spacing=-360, line_spacing=1.0)
    add_text(slide, "Proaktiv. Persönlich.",
             Inches(0), Inches(3.1), SLIDE_W, Inches(1.4),
             font=FONT_SERIF, size_pt=66, color=TEAL,
             italic=True, align="center", letter_spacing=-360, line_spacing=1.0)

    # Subline (centered, max-width)
    add_text(slide,
             "iOS-App für Jobsuchende & Schüler — gebaut für den User, nicht für Unternehmen.",
             Inches(1.5), Inches(4.55), SLIDE_W - Inches(3), Inches(0.8),
             font=FONT_SANS, size_pt=16, color=TEXT2, align="center",
             line_spacing=1.4)

    # Logo + Wordmark bottom-left
    if ASSET_LOGO.exists():
        slide.shapes.add_picture(str(ASSET_LOGO),
                                 MARGIN_X, SLIDE_H - Inches(0.95),
                                 height=Inches(0.45))
    add_text(slide, "ApplyOS",
             MARGIN_X + Inches(0.6), SLIDE_H - Inches(0.88),
             Inches(3), Inches(0.4),
             font=FONT_SERIF, size_pt=17, color=TEXT,
             italic=True, bold=True, letter_spacing=-20)

    # Brand line bottom-right
    add_text(slide, "tryapplyos.app · Vertraulich",
             SLIDE_W - MARGIN_X - Inches(4),
             SLIDE_H - Inches(0.75), Inches(4), Inches(0.3),
             font=FONT_SANS, size_pt=9, color=FAINT,
             align="right", letter_spacing=160)


def build_slide_02_problem(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Das Problem", MARGIN_X, MARGIN_TOP)
    add_serif_heading(
        slide,
        "Karriereportale & Berufsorientierung sind für die falschen Kunden gebaut.",
        MARGIN_X, MARGIN_TOP + Inches(0.3),
        SLIDE_W - 2 * MARGIN_X, Inches(2.0),
        size_pt=34, color=TEXT, letter_spacing=-240,
    )

    # Dual-Column: links Jobsuchende (Teal), rechts Schüler (Tan)
    col_y = Inches(3.2)
    col_h = Inches(3.4)
    col_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.4)) // 2

    # Linke Spalte
    left_x = MARGIN_X
    add_card(slide, left_x, col_y, col_w, col_h,
             fill_color=PAPER, border_color=HAIRLINE_CARD, radius_pt=14)
    add_eyebrow(slide, "Jobsuchende", left_x + Inches(0.4),
                col_y + Inches(0.35), color=TEAL, width=col_w - Inches(0.8))

    left_bullets_top = col_y + Inches(0.85)
    bullets_jobs = [
        ("Optimiert für Unternehmen.",
         "User werden als Ressource vermittelt — Profil zweitrangig."),
        ("Kognitive Last beim User.",
         "Keywords, Filter, manuelles Tracking — der User pflegt das System."),
        ("Vermittler verdienen, nicht der User.",
         "Plattformen monetarisieren Recruiter-Zugang, nicht User-Erfolg."),
    ]
    bullet_y = left_bullets_top
    for headline, body in bullets_jobs:
        add_text(slide, "—", left_x + Inches(0.4), bullet_y,
                 Inches(0.4), Inches(0.4),
                 font=FONT_SERIF, size_pt=18, color=TEAL)
        add_text(slide, headline, left_x + Inches(0.7), bullet_y,
                 col_w - Inches(1.1), Inches(0.4),
                 font=FONT_SANS, size_pt=13, color=TEXT, bold=True,
                 letter_spacing=-10)
        add_text(slide, body, left_x + Inches(0.7), bullet_y + Inches(0.32),
                 col_w - Inches(1.1), Inches(0.5),
                 font=FONT_SANS, size_pt=11, color=TEXT2, line_spacing=1.35)
        bullet_y += Inches(0.78)

    # Rechte Spalte
    right_x = MARGIN_X + col_w + Inches(0.4)
    add_card(slide, right_x, col_y, col_w, col_h,
             fill_color=PAPER, border_color=HAIRLINE_CARD, radius_pt=14)
    add_eyebrow(slide, "Schüler", right_x + Inches(0.4),
                col_y + Inches(0.35), color=TAN, width=col_w - Inches(0.8))

    bullets_schueler = [
        ("90er-Jahre-Multiple-Choice.",
         "Berufstests fragen nach Lieblingsfächern — nicht nach Werten & Stärken."),
        ("Beratung mit Vermittlungsinteresse.",
         "Karriere-Beratungen verdienen an Vermittlung — nicht an Klarheit."),
        ("Eltern- & Schul-Druck statt Selbsterkenntnis.",
         '„Was würde sich gut machen?" statt „Was passt zu mir?"'),
    ]
    bullet_y = left_bullets_top
    for headline, body in bullets_schueler:
        add_text(slide, "—", right_x + Inches(0.4), bullet_y,
                 Inches(0.4), Inches(0.4),
                 font=FONT_SERIF, size_pt=18, color=TAN)
        add_text(slide, headline, right_x + Inches(0.7), bullet_y,
                 col_w - Inches(1.1), Inches(0.4),
                 font=FONT_SANS, size_pt=13, color=TEXT, bold=True,
                 letter_spacing=-10)
        add_text(slide, body, right_x + Inches(0.7), bullet_y + Inches(0.32),
                 col_w - Inches(1.1), Inches(0.5),
                 font=FONT_SANS, size_pt=11, color=TEXT2, line_spacing=1.35)
        bullet_y += Inches(0.78)

    add_footer(slide, 2)


def build_slide_03_solution(prs):
    slide = new_slide(prs)
    add_glow(slide, SLIDE_W // 2, Inches(3.0),
             Inches(11), Inches(5.5), color=TEAL_GLOW)
    add_eyebrow(slide, "Die Lösung", MARGIN_X, MARGIN_TOP)

    add_serif_heading_dual(
        slide, "Ein KI-Agent.", "Zwei Lebensphasen.",
        MARGIN_X, Inches(1.3), SLIDE_W - 2 * MARGIN_X,
        size_pt=68, color1=TEXT, color2=TEAL, letter_spacing=-380,
    )

    add_text(
        slide,
        "ApplyOS arbeitet täglich für dich — egal ob du den nächsten Job suchst oder "
        "entscheidest, was du nach der Schule machst.",
        MARGIN_X, Inches(4.0), SLIDE_W - 2 * MARGIN_X, Inches(1.0),
        font=FONT_SANS, size_pt=18, color=TEXT2, line_spacing=1.45,
    )

    # Mini-Stats Row
    stats_y = Inches(5.3)
    stats = [
        ("100%", "USER-FOKUS"),
        ("0", "UNTERNEHMENSKUNDEN"),
        ("↓ 85%", "AUFWAND"),
    ]
    col_w_stat = (SLIDE_W - 2 * MARGIN_X - Inches(0.8)) // 3
    for i, (num, label) in enumerate(stats):
        x = MARGIN_X + i * (col_w_stat + Inches(0.4))
        add_card(slide, x, stats_y, col_w_stat, Inches(1.5),
                 fill_color=PAPER, border_color=HAIRLINE_CARD, radius_pt=12)
        add_text(slide, num, x, stats_y + Inches(0.2),
                 col_w_stat, Inches(0.8),
                 font=FONT_SERIF, size_pt=42, color=TEAL,
                 italic=True, align="center", letter_spacing=-160)
        add_text(slide, label, x, stats_y + Inches(1.05),
                 col_w_stat, Inches(0.3),
                 font=FONT_SANS, size_pt=9, color=MUTED,
                 align="center", letter_spacing=220)

    add_footer(slide, 3)


def build_slide_04_market(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Warum jetzt", MARGIN_X, MARGIN_TOP)
    add_serif_heading(
        slide, "Drei Märkte. Ein Moment.",
        MARGIN_X, MARGIN_TOP + Inches(0.3),
        SLIDE_W - 2 * MARGIN_X, Inches(1.0),
        size_pt=44, color=TEXT, letter_spacing=-280,
    )

    cards = [
        ("01", "GenAI macht Agent-UX möglich.",
         "Large Language Models sind consumer-grade, mobile-fähig und in der "
         "Lage, proaktiv für den User zu arbeiten — nicht nur zu reagieren."),
        ("02", "Application Fatigue auf Rekord.",
         "Multi-Kanal-Bewerbungen, ATS-Filter, ghost-Stellen. User sind erschöpft "
         "und suchen nach Tools, die für sie filtern, nicht gegen sie."),
        ("03", "Zwei Märkte — eine Lösung.",
         "DE-Arbeitsmarkt mit hunderttausenden offenen Stellen täglich plus "
         "~700.000 Schulabgänger jährlich — keine App adressiert beide."),
    ]
    card_y = Inches(3.0)
    card_h = Inches(3.4)
    card_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.6)) // 3

    for i, (num, headline, body) in enumerate(cards):
        x = MARGIN_X + i * (card_w + Inches(0.3))
        add_card(slide, x, card_y, card_w, card_h,
                 fill_color=PAPER, border_color=HAIRLINE_CARD, radius_pt=14)
        # Number
        add_text(slide, num, x + Inches(0.4), card_y + Inches(0.35),
                 Inches(1), Inches(0.7),
                 font=FONT_SERIF, size_pt=36, color=TEAL,
                 italic=True, letter_spacing=-160)
        # Headline
        add_text(slide, headline, x + Inches(0.4), card_y + Inches(1.15),
                 card_w - Inches(0.8), Inches(0.95),
                 font=FONT_SERIF, size_pt=20, color=TEXT,
                 letter_spacing=-80, line_spacing=1.1)
        # Body
        add_text(slide, body, x + Inches(0.4), card_y + Inches(2.15),
                 card_w - Inches(0.8), Inches(1.2),
                 font=FONT_SANS, size_pt=11, color=TEXT2, line_spacing=1.45)

    # Footnote
    add_text(slide,
             "Stellen aus einer der größten Jobdatenbanken Deutschlands (Bundesagentur-API).",
             MARGIN_X, Inches(6.55), SLIDE_W - 2 * MARGIN_X, Inches(0.3),
             font=FONT_SANS, size_pt=9, color=FAINT,
             italic=True, letter_spacing=80)

    add_footer(slide, 4)


# ──────────────────────────────────────────────────────────────────────────────
# Phone Mockup Composers (für Slide 5)
# ──────────────────────────────────────────────────────────────────────────────

APP_TEAL = RGBColor(0x0D, 0x9C, 0x81)
APP_TEAL_SOFT = RGBColor(0xDB, 0xEE, 0xE4)
APP_BG = RGBColor(0xF4, 0xF3, 0xEF)
APP_TEXT = RGBColor(0x0A, 0x0A, 0x0A)
APP_TEXT2 = RGBColor(0x3A, 0x3A, 0x3A)
APP_DIM = RGBColor(0x88, 0x88, 0x88)
APP_DIVIDER = RGBColor(0xF1, 0xEF, 0xE7)


def draw_phone_radar(slide, sl, st, sw, sh):
    """Zeichnet den Radar/Jobsuche-Screen in die Phone-Bounds."""
    pad = Pt(4)
    inner_x = sl + pad
    inner_y = st + pad
    inner_w = sw - 2 * pad
    # Header
    add_text(slide, "Radar", inner_x + Pt(8), inner_y + Pt(10),
             inner_w - Pt(16), Pt(20),
             font=FONT_SANS, size_pt=10, color=APP_TEXT, bold=True)
    add_text(slide, "Heute 12 neue", inner_x + Pt(8), inner_y + Pt(24),
             inner_w - Pt(16), Pt(14),
             font=FONT_SANS, size_pt=6, color=APP_DIM)
    # 3 Job-Cards
    card_top = inner_y + Pt(46)
    for i in range(3):
        cy = card_top + i * Pt(48)
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   inner_x + Pt(6), cy,
                                   inner_w - Pt(12), Pt(42))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        c.line.color.rgb = APP_DIVIDER
        c.line.width = Pt(0.5)
        c.adjustments[0] = 0.18
        remove_shadow(c)
        # Title
        titles = ["Senior iOS Engineer", "Product Designer", "Data Scientist"]
        companies = ["TechCo · Berlin", "DesignStudio · Remote", "FinTech · München"]
        scores = ["92%", "87%", "81%"]
        add_text(slide, titles[i], inner_x + Pt(12), cy + Pt(6),
                 inner_w - Pt(56), Pt(14),
                 font=FONT_SANS, size_pt=7, color=APP_TEXT, bold=True)
        add_text(slide, companies[i], inner_x + Pt(12), cy + Pt(20),
                 inner_w - Pt(56), Pt(14),
                 font=FONT_SANS, size_pt=5.5, color=APP_DIM)
        # Score badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       inner_x + inner_w - Pt(38),
                                       cy + Pt(8), Pt(28), Pt(20))
        badge.fill.solid()
        badge.fill.fore_color.rgb = APP_TEAL_SOFT
        badge.line.fill.background()
        badge.adjustments[0] = 0.35
        remove_shadow(badge)
        add_text(slide, scores[i],
                 inner_x + inner_w - Pt(38), cy + Pt(10),
                 Pt(28), Pt(16),
                 font=FONT_SANS, size_pt=6, color=APP_TEAL, bold=True,
                 align="center")


def draw_phone_job_detail(slide, sl, st, sw, sh):
    """Job-Detail-Screen mit Match-Score-Visualisierung."""
    pad = Pt(4)
    inner_x = sl + pad
    inner_y = st + pad
    inner_w = sw - 2 * pad
    # Title
    add_text(slide, "Senior iOS Engineer",
             inner_x + Pt(10), inner_y + Pt(12),
             inner_w - Pt(20), Pt(18),
             font=FONT_SANS, size_pt=8, color=APP_TEXT, bold=True)
    add_text(slide, "TechCo · Berlin", inner_x + Pt(10), inner_y + Pt(26),
             inner_w - Pt(20), Pt(12),
             font=FONT_SANS, size_pt=6, color=APP_DIM)
    # Big Score circle
    cx = inner_x + inner_w // 2
    cy = inner_y + Pt(70)
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  cx - Pt(28), cy - Pt(28),
                                  Pt(56), Pt(56))
    ring.fill.solid()
    ring.fill.fore_color.rgb = APP_TEAL_SOFT
    ring.line.color.rgb = APP_TEAL
    ring.line.width = Pt(2.5)
    remove_shadow(ring)
    add_text(slide, "87%", cx - Pt(28), cy - Pt(14), Pt(56), Pt(28),
             font=FONT_SERIF, size_pt=18, color=APP_TEAL, bold=True,
             italic=True, align="center")
    add_text(slide, "Match-Score", inner_x, cy + Pt(32), inner_w, Pt(14),
             font=FONT_SANS, size_pt=6, color=APP_DIM, align="center",
             letter_spacing=120)
    # Mini criteria bars
    crit_top = cy + Pt(54)
    criteria = [("Fachrichtung", 0.95), ("Standort", 1.0), ("Skills", 0.78)]
    for i, (label, fill) in enumerate(criteria):
        ry = crit_top + i * Pt(18)
        add_text(slide, label, inner_x + Pt(10), ry,
                 Pt(50), Pt(12),
                 font=FONT_SANS, size_pt=5.5, color=APP_TEXT2)
        bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         inner_x + Pt(60), ry + Pt(3),
                                         inner_w - Pt(80), Pt(5))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = APP_DIVIDER
        bar_bg.line.fill.background()
        bar_bg.adjustments[0] = 0.5
        remove_shadow(bar_bg)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      inner_x + Pt(60), ry + Pt(3),
                                      int((inner_w - Pt(80)) * fill), Pt(5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = APP_TEAL
        bar.line.fill.background()
        bar.adjustments[0] = 0.5
        remove_shadow(bar)


def draw_phone_pipeline(slide, sl, st, sw, sh):
    """Pipeline-Screen mit Kanban-Spalten."""
    pad = Pt(4)
    inner_x = sl + pad
    inner_y = st + pad
    inner_w = sw - 2 * pad
    # Header
    add_text(slide, "Pipeline", inner_x + Pt(8), inner_y + Pt(10),
             inner_w - Pt(16), Pt(18),
             font=FONT_SANS, size_pt=10, color=APP_TEXT, bold=True)
    add_text(slide, "4 aktiv", inner_x + Pt(8), inner_y + Pt(24),
             inner_w - Pt(16), Pt(14),
             font=FONT_SANS, size_pt=6, color=APP_DIM)
    # 4 Spalten (vertikal gestapelt für Phone-View)
    stages = [
        ("ENTDECKT", 12, APP_DIM),
        ("BEWORBEN", 5, APP_TEAL),
        ("INTERVIEW", 2, RGBColor(0xD7, 0x7B, 0x1C)),
        ("ANGEBOT", 1, RGBColor(0x5A, 0x6C, 0xFF)),
    ]
    stage_y = inner_y + Pt(46)
    for i, (label, count, color) in enumerate(stages):
        sy = stage_y + i * Pt(34)
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   inner_x + Pt(6), sy,
                                   inner_w - Pt(12), Pt(28))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        c.line.color.rgb = APP_DIVIDER
        c.line.width = Pt(0.5)
        c.adjustments[0] = 0.25
        remove_shadow(c)
        # Color dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     inner_x + Pt(12), sy + Pt(10),
                                     Pt(8), Pt(8))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        remove_shadow(dot)
        add_text(slide, label, inner_x + Pt(24), sy + Pt(6),
                 Pt(60), Pt(14),
                 font=FONT_SANS, size_pt=5.5, color=APP_TEXT2,
                 bold=True, letter_spacing=120)
        add_text(slide, str(count),
                 inner_x + inner_w - Pt(28), sy + Pt(5),
                 Pt(20), Pt(18),
                 font=FONT_SERIF, size_pt=12, color=APP_TEXT,
                 italic=True, align="right")


def build_slide_05_product_jobsuchende(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Produkt — Jobsuchende", MARGIN_X, MARGIN_TOP)
    add_serif_heading(
        slide, "Die App, die für dich arbeitet.",
        MARGIN_X, MARGIN_TOP + Inches(0.3),
        SLIDE_W - 2 * MARGIN_X, Inches(0.9),
        size_pt=38, color=TEXT, letter_spacing=-240,
    )

    cards = [
        ("Proaktive Jobsuche",
         "Täglich frische Top-Matches, automatisch auf dich zugeschnitten — keine "
         "Filter, keine Keywords.",
         draw_phone_radar),
        ("Intelligentes Matching",
         "Ein KI-Score zeigt auf einen Blick, wie gut jede Stelle wirklich zu "
         "deinem Profil passt.",
         draw_phone_job_detail),
        ("Kanban-Pipeline",
         "Alle Bewerbungen auf einen Blick: Entdeckt, Beworben, Interview, "
         "Angebot.",
         draw_phone_pipeline),
    ]
    card_y = Inches(2.5)
    card_h = Inches(4.3)
    card_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.5)) // 3
    phone_w = Inches(1.6)
    phone_h = Inches(2.6)

    for i, (title, body, draw_phone) in enumerate(cards):
        x = MARGIN_X + i * (card_w + Inches(0.25))
        add_card(slide, x, card_y, card_w, card_h,
                 fill_color=PAPER, border_color=HAIRLINE_CARD, radius_pt=14)
        # Phone centered horizontally in card
        phone_x = x + (card_w - phone_w) // 2
        phone_y = card_y + Inches(0.4)
        _, _, screen_bounds = add_phone_frame(slide, phone_x, phone_y, phone_w, phone_h)
        # Draw screen content
        draw_phone(slide, *screen_bounds)
        # Divider
        add_hairline(slide, x + Inches(0.4), card_y + Inches(3.15),
                     card_w - Inches(0.8), color=HAIRLINE_CARD)
        # Title
        add_text(slide, title, x + Inches(0.4), card_y + Inches(3.35),
                 card_w - Inches(0.8), Inches(0.4),
                 font=FONT_SERIF, size_pt=18, color=TEXT,
                 letter_spacing=-60)
        # Body
        add_text(slide, body, x + Inches(0.4), card_y + Inches(3.78),
                 card_w - Inches(0.8), Inches(0.7),
                 font=FONT_SANS, size_pt=10.5, color=TEXT2, line_spacing=1.45)

    add_footer(slide, 5)


def build_slide_06_product_schueler(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Produkt — Schüler", MARGIN_X, MARGIN_TOP, color=TAN)
    add_serif_heading(
        slide, "Wissenschaftlich fundiertes Profiling. In 5 Minuten.",
        MARGIN_X, MARGIN_TOP + Inches(0.3),
        SLIDE_W - 2 * MARGIN_X, Inches(0.9),
        size_pt=32, color=TEXT, letter_spacing=-220,
    )

    # Linke Hälfte: 4-Phasen-Flow (vertikal stacked)
    flow_left = MARGIN_X
    flow_top = Inches(2.4)
    flow_w = Inches(7.6)

    phases = [
        ("01", "Hook",
         "3 Swipe-Fragen",
         'z.B. „Wie sieht dein Schreibtisch aus?" — misst Lernstil & Format-Präferenz.',
         TAN),
        ("02", "Aha Moment",
         "Archetyp-Reveal",
         '„Du bist ein Tüftler!" — erste Selbsterkenntnis nach 90 Sekunden.',
         TAN),
        ("03", "Deep Dive",
         "RIASEC + Status-quo",
         "3 Fragen zu Holland-Dimensionen + 1 Frage zur Umfeld-Erwartung.",
         TAN),
        ("04", "Synthese",
         "Standort + Berechnung + Ergebnis",
         "Output: RIASEC-Profil-Label, Format-Empfehlung, Top-Fachrichtungen.",
         TAN),
    ]

    row_h = Inches(0.9)
    for i, (num, phase, sub, body, accent) in enumerate(phases):
        y = flow_top + i * row_h
        # Number
        add_text(slide, num, flow_left, y + Inches(0.05),
                 Inches(0.7), Inches(0.7),
                 font=FONT_SERIF, size_pt=30, color=accent,
                 italic=True, letter_spacing=-100)
        # Phase title
        add_text(slide, phase, flow_left + Inches(0.8), y,
                 Inches(2.0), Inches(0.4),
                 font=FONT_SERIF, size_pt=18, color=TEXT,
                 letter_spacing=-60)
        # Sub-label
        add_text(slide, sub, flow_left + Inches(0.8), y + Inches(0.38),
                 Inches(2.4), Inches(0.3),
                 font=FONT_SANS, size_pt=9, color=accent,
                 bold=True, letter_spacing=140)
        # Body
        add_text(slide, body, flow_left + Inches(3.3), y + Inches(0.1),
                 flow_w - Inches(3.3), Inches(0.7),
                 font=FONT_SANS, size_pt=11, color=TEXT2, line_spacing=1.4)
        # Hairline between phases
        if i < len(phases) - 1:
            add_hairline(slide, flow_left + Inches(0.8), y + Inches(0.82),
                         flow_w - Inches(0.8), color=HAIRLINE_BG)

    # Footnote
    add_text(slide,
             "Basiert auf Holland's RIASEC-Theorie — wissenschaftlich etabliert, "
             "weltweit in Karriereberatung genutzt.",
             flow_left, Inches(6.3), flow_w, Inches(0.3),
             font=FONT_SANS, size_pt=9, color=FAINT,
             italic=True, letter_spacing=80)

    # Rechte Hälfte: app-result.png als App-Screen
    image_x = MARGIN_X + flow_w + Inches(0.3)
    image_w = SLIDE_W - MARGIN_X - image_x
    image_y = Inches(2.2)
    image_h = Inches(4.4)

    if ASSET_RESULT.exists():
        # Card-Background hinter dem Bild
        card_bg = add_card(slide, image_x, image_y, image_w, image_h,
                           fill_color=PAPER, border_color=HAIRLINE_CARD,
                           radius_pt=14)
        # Bild zentriert mit Padding (max width passend)
        pad = Pt(10)
        img_max_w = image_w - 2 * pad
        img_max_h = image_h - 2 * pad
        # Lade Bild zum Ratio-Check
        from PIL import Image as PILImage
        try:
            with PILImage.open(ASSET_RESULT) as im:
                w_px, h_px = im.size
                ratio = w_px / h_px
        except Exception:
            ratio = 0.5  # iPhone-Schätzung
        # Skaliere passend
        if img_max_w / img_max_h > ratio:
            # Höhe limitiert
            new_h = img_max_h
            new_w = int(new_h * ratio)
        else:
            new_w = img_max_w
            new_h = int(new_w / ratio)
        img_left = image_x + (image_w - new_w) // 2
        img_top = image_y + (image_h - new_h) // 2
        slide.shapes.add_picture(str(ASSET_RESULT),
                                 img_left, img_top, width=new_w, height=new_h)

        # Caption unter dem Bild
        add_text(slide, "Beispiel-Output der Synthese-Phase",
                 image_x, Inches(6.65), image_w, Inches(0.3),
                 font=FONT_SANS, size_pt=9, color=FAINT,
                 italic=True, align="center", letter_spacing=80)

    add_footer(slide, 6)


def build_slide_07_differentiation(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Warum ApplyOS", MARGIN_X, MARGIN_TOP)
    add_serif_heading_dual(
        slide, "Weil du es", "wert bist.",
        MARGIN_X, Inches(1.2), SLIDE_W - 2 * MARGIN_X,
        size_pt=58, color1=TEXT, color2=TEAL, letter_spacing=-360,
    )

    # 2x2 USP-Grid
    grid_top = Inches(3.4)
    grid_w = SLIDE_W - 2 * MARGIN_X
    cell_w = (grid_w - Inches(0.3)) // 2
    cell_h = Inches(1.55)

    usps = [
        ("100% User-Fokus.",
         "Keine Unternehmens- oder Beratungskunden. Kein Vermittlungsinteresse. "
         "Die App verdient nicht an dir — sie arbeitet für dich.",
         TEAL),
        ("Privacy by Design.",
         "Deine Daten bleiben bei dir. ApplyOS leitet dich zur Original-Anzeige "
         "weiter — keine Recruiter-Durchleitung, kein Profil-Verkauf.",
         TEAL),
        ("Wissenschaftlich fundiert.",
         "RIASEC-Profiling für Schüler. 5-Kriterien-Match-Score: Fachrichtung 30% · "
         "Standort 25% · Keywords 20% · Startdatum 15% · Aktualität 10%.",
         TEAL),
        ("Dual-Audience aus einem Stack.",
         "Jobsuchende und Schüler aus derselben SwiftUI-Codebase und geteilten "
         "AI-Infra — zwei Märkte, ein Produkt, eine Engineering-Investition.",
         TEAL),
    ]
    for i, (title, body, accent) in enumerate(usps):
        col = i % 2
        row = i // 2
        x = MARGIN_X + col * (cell_w + Inches(0.3))
        y = grid_top + row * (cell_h + Inches(0.3))
        add_card(slide, x, y, cell_w, cell_h,
                 fill_color=PAPER, border_color=HAIRLINE_CARD, radius_pt=14)
        # Accent line links
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            x + Inches(0.4), y + Inches(0.32),
                                            Inches(0.04), Inches(0.9))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()
        remove_shadow(accent_bar)
        # Title
        add_text(slide, title, x + Inches(0.6), y + Inches(0.3),
                 cell_w - Inches(0.9), Inches(0.5),
                 font=FONT_SERIF, size_pt=20, color=TEXT,
                 letter_spacing=-80)
        # Body
        add_text(slide, body, x + Inches(0.6), y + Inches(0.78),
                 cell_w - Inches(0.9), cell_h - Inches(0.9),
                 font=FONT_SANS, size_pt=11, color=TEXT2, line_spacing=1.45)

    add_footer(slide, 7)


def build_slide_08_gtm(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Go-to-Market & Business Model", MARGIN_X, MARGIN_TOP)
    add_serif_heading(
        slide, "Closed Beta heute — App Store morgen.",
        MARGIN_X, MARGIN_TOP + Inches(0.3),
        SLIDE_W - 2 * MARGIN_X, Inches(0.9),
        size_pt=36, color=TEXT, letter_spacing=-240,
    )

    col_y = Inches(2.7)
    col_h = Inches(4.0)
    col_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.4)) // 2

    # Linke Spalte: Distribution
    left_x = MARGIN_X
    add_card(slide, left_x, col_y, col_w, col_h,
             fill_color=PAPER, border_color=HAIRLINE_CARD, radius_pt=14)
    add_eyebrow(slide, "Distribution", left_x + Inches(0.4),
                col_y + Inches(0.35), color=TEAL)
    add_text(slide, "iOS-First.",
             left_x + Inches(0.4), col_y + Inches(0.85),
             col_w - Inches(0.8), Inches(0.5),
             font=FONT_SERIF, size_pt=22, color=TEXT, letter_spacing=-80)
    bullets_dist = [
        "SwiftUI-native — performant, schnell iteriert.",
        "App Store als zentraler Distributions-Kanal.",
        "Organische Akquise via tryapplyos.app.",
        "SEO auf DE-Karriere-Keywords (Jobsuche, Berufsorientierung).",
        "Closed Beta in Wellen — Hype + Feedback-Loop.",
        "EN-Markt parallel veröffentlicht — internationale Reichweite.",
    ]
    bullet_y = col_y + Inches(1.55)
    for b in bullets_dist:
        add_text(slide, "·", left_x + Inches(0.4), bullet_y,
                 Inches(0.2), Inches(0.3),
                 font=FONT_SANS, size_pt=14, color=TEAL, bold=True)
        add_text(slide, b, left_x + Inches(0.6), bullet_y,
                 col_w - Inches(1.0), Inches(0.32),
                 font=FONT_SANS, size_pt=11, color=TEXT2, line_spacing=1.4)
        bullet_y += Inches(0.36)

    # Rechte Spalte: Monetization
    right_x = MARGIN_X + col_w + Inches(0.4)
    add_card(slide, right_x, col_y, col_w, col_h,
             fill_color=PAPER, border_color=HAIRLINE_CARD, radius_pt=14)
    add_eyebrow(slide, "Monetization", right_x + Inches(0.4),
                col_y + Inches(0.35), color=TEAL)
    add_text(slide, "Free + Premium.",
             right_x + Inches(0.4), col_y + Inches(0.85),
             col_w - Inches(0.8), Inches(0.5),
             font=FONT_SERIF, size_pt=22, color=TEXT, letter_spacing=-80)
    bullets_mon = [
        "Free-Tier: Profiling, Matching, Pipeline — alle Kern-Features.",
        "Premium-Tier in Evaluation: erweiterte AI-Generierung.",
        "Coaching-Features als künftiger Upsell (Interview-Coach).",
        "Pricing wird mit Beta-Feedback finalisiert.",
        "B2C-Modell — keine B2B-Vermittlung, keine Recruiter-Provisionen.",
        "Skalierbar via geringe Marginal-Kosten (Hosting + AI-API).",
    ]
    bullet_y = col_y + Inches(1.55)
    for b in bullets_mon:
        add_text(slide, "·", right_x + Inches(0.4), bullet_y,
                 Inches(0.2), Inches(0.3),
                 font=FONT_SANS, size_pt=14, color=TEAL, bold=True)
        add_text(slide, b, right_x + Inches(0.6), bullet_y,
                 col_w - Inches(1.0), Inches(0.32),
                 font=FONT_SANS, size_pt=11, color=TEXT2, line_spacing=1.4)
        bullet_y += Inches(0.36)

    add_footer(slide, 8)


def draw_phone_waitlist(slide, sl, st, sw, sh):
    """Waitlist-Confirmation-Screen."""
    pad = Pt(4)
    inner_x = sl + pad
    inner_y = st + pad
    inner_w = sw - 2 * pad
    inner_h = sh - 2 * pad
    cx = inner_x + inner_w // 2
    # Checkmark badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   cx - Pt(24), inner_y + Pt(36),
                                   Pt(48), Pt(48))
    badge.fill.solid()
    badge.fill.fore_color.rgb = APP_TEAL_SOFT
    badge.line.color.rgb = APP_TEAL
    badge.line.width = Pt(2)
    remove_shadow(badge)
    add_text(slide, "✓", cx - Pt(24), inner_y + Pt(46), Pt(48), Pt(28),
             font=FONT_SANS, size_pt=24, color=APP_TEAL, bold=True,
             align="center")
    # Headline
    add_text(slide, "Willkommen", inner_x, inner_y + Pt(98),
             inner_w, Pt(20),
             font=FONT_SERIF, size_pt=12, color=APP_TEXT, align="center")
    add_text(slide, "an Bord.", inner_x, inner_y + Pt(112),
             inner_w, Pt(20),
             font=FONT_SERIF, size_pt=12, color=APP_TEAL, italic=True,
             align="center")
    # Position
    add_text(slide, "DEINE POSITION", inner_x, inner_y + Pt(150),
             inner_w, Pt(14),
             font=FONT_SANS, size_pt=5.5, color=APP_DIM,
             align="center", letter_spacing=200)
    add_text(slide, "Platz 42 von 100", inner_x, inner_y + Pt(164),
             inner_w, Pt(20),
             font=FONT_SERIF, size_pt=14, color=APP_TEXT, bold=True,
             italic=True, align="center")
    # CTA hint
    add_text(slide, "Beta-Start: in 24h", inner_x, inner_y + Pt(196),
             inner_w, Pt(14),
             font=FONT_SANS, size_pt=6, color=APP_DIM,
             align="center", letter_spacing=100)


def build_slide_09_traction(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Traction", MARGIN_X, MARGIN_TOP)
    add_serif_heading(
        slide, "Closed Beta startet.",
        MARGIN_X, MARGIN_TOP + Inches(0.3),
        SLIDE_W - 2 * MARGIN_X, Inches(0.9),
        size_pt=40, color=TEXT, letter_spacing=-260,
    )

    # Linke Hälfte: Big Stat + Bullets
    left_x = MARGIN_X
    left_w = Inches(8.0)

    # Big Stat
    add_text(slide, "100", left_x, Inches(2.7), left_w, Inches(1.8),
             font=FONT_SERIF, size_pt=140, color=TEAL,
             italic=True, letter_spacing=-540, line_spacing=1.0)
    add_text(slide, "Exklusive Beta-Plätze",
             left_x + Inches(2.6), Inches(3.5),
             Inches(5), Inches(0.4),
             font=FONT_SANS, size_pt=11, color=MUTED,
             letter_spacing=200)

    # Bullets
    bullets = [
        ("Waitlist live", "tryapplyos.app mit Live-Counter & Wellen-Onboarding."),
        ("Direkter Draht zum Team", "Beta-User formen das Produkt mit."),
        ("DE-First, EN parallel", "Internationale Reichweite ab Tag eins."),
        ("App Store-Submission läuft", "Public Launch unmittelbar nach Beta."),
    ]
    bullet_y = Inches(4.65)
    for headline, body in bullets:
        add_text(slide, "—", left_x, bullet_y, Inches(0.3), Inches(0.3),
                 font=FONT_SERIF, size_pt=14, color=TEAL)
        add_text(slide, headline, left_x + Inches(0.3), bullet_y,
                 Inches(2.5), Inches(0.3),
                 font=FONT_SANS, size_pt=11, color=TEXT, bold=True)
        add_text(slide, body, left_x + Inches(2.9), bullet_y,
                 left_w - Inches(2.9), Inches(0.3),
                 font=FONT_SANS, size_pt=11, color=TEXT2)
        bullet_y += Inches(0.4)

    # Rechte Hälfte: Mini-Phone mit Waitlist-Screen
    phone_w = Inches(2.4)
    phone_h = Inches(4.0)
    phone_x = SLIDE_W - MARGIN_X - phone_w - Inches(0.8)
    phone_y = Inches(2.5)
    _, _, screen_bounds = add_phone_frame(slide, phone_x, phone_y, phone_w, phone_h)
    draw_phone_waitlist(slide, *screen_bounds)

    add_footer(slide, 9)


def build_slide_10_roadmap(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Roadmap", MARGIN_X, MARGIN_TOP)
    add_serif_heading(
        slide, "Was als nächstes kommt.",
        MARGIN_X, MARGIN_TOP + Inches(0.3),
        SLIDE_W - 2 * MARGIN_X, Inches(0.9),
        size_pt=40, color=TEXT, letter_spacing=-260,
    )

    # Horizontale Timeline
    milestones = [
        ("Q2 2026", "Closed Beta",
         "Erste 100 User. Feedback-Iteration. Onboarding-Optimierung.", TEAL, True),
        ("Q3 2026", "Public iOS Launch",
         "DE + EN im App Store. Organic Growth via SEO + Press.", TEAL, False),
        ("Q4 2026", "Premium-Tier",
         "Erweiterte AI-Features: Anschreiben-Personalisierung, Interview-Coach.", TEAL, False),
        ("2027", "Android & Expansion",
         "Native Android-App. Internationalisierung in weitere Märkte.", TEAL, False),
    ]
    timeline_y = Inches(3.4)
    timeline_left = MARGIN_X + Inches(0.4)
    timeline_right = SLIDE_W - MARGIN_X - Inches(0.4)
    timeline_w = timeline_right - timeline_left

    # Horizontal line
    line_y = timeline_y + Inches(0.45)
    add_hairline(slide, timeline_left, line_y, timeline_w, color=HAIRLINE_BG,
                 width_pt=1.0)

    col_w = timeline_w // len(milestones)
    for i, (quarter, headline, body, accent, current) in enumerate(milestones):
        cx = timeline_left + i * col_w + col_w // 2
        # Dot
        dot_size = Pt(14 if current else 10)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     cx - dot_size // 2, line_y - dot_size // 2,
                                     dot_size, dot_size)
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent if current else PAPER
        dot.line.color.rgb = accent
        dot.line.width = Pt(1.5)
        remove_shadow(dot)
        # Quarter above
        add_text(slide, quarter, cx - col_w // 2, timeline_y - Inches(0.05),
                 col_w, Inches(0.4),
                 font=FONT_SANS, size_pt=10, color=accent, bold=True,
                 align="center", letter_spacing=200)
        # Headline below
        add_text(slide, headline, cx - col_w // 2, line_y + Inches(0.3),
                 col_w, Inches(0.5),
                 font=FONT_SERIF, size_pt=18, color=TEXT,
                 align="center", letter_spacing=-60)
        # Body below
        add_text(slide, body, cx - col_w // 2 + Inches(0.2),
                 line_y + Inches(0.8),
                 col_w - Inches(0.4), Inches(1.5),
                 font=FONT_SANS, size_pt=10, color=TEXT2,
                 align="center", line_spacing=1.4)

    add_footer(slide, 10)


def build_slide_11_team(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Team", MARGIN_X, MARGIN_TOP)
    add_serif_heading(
        slide, "Wer ApplyOS baut.",
        MARGIN_X, MARGIN_TOP + Inches(0.3),
        SLIDE_W - 2 * MARGIN_X, Inches(0.9),
        size_pt=40, color=TEXT, letter_spacing=-260,
    )

    members = [
        ("Florian Krause", "Founder · Product & Engineering",
         "Bringt iOS-Engineering und Produkt-Vision zusammen. Gebaut, weil "
         "klassische Karriereportale ihn selbst gestört haben."),
        ("[Name]", "[Rolle]",
         "[Bio — Florian füllt manuell ein.]"),
        ("[Name]", "[Rolle]",
         "[Bio — Florian füllt manuell ein.]"),
    ]
    card_y = Inches(2.7)
    card_h = Inches(4.0)
    card_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.6)) // 3

    for i, (name, role, bio) in enumerate(members):
        x = MARGIN_X + i * (card_w + Inches(0.3))
        add_card(slide, x, card_y, card_w, card_h,
                 fill_color=PAPER, border_color=HAIRLINE_CARD, radius_pt=14)
        # Avatar-Platzhalter
        ax = x + (card_w - Inches(1.6)) // 2
        avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        ax, card_y + Inches(0.5),
                                        Inches(1.6), Inches(1.6))
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = PAPER_HI
        avatar.line.color.rgb = HAIRLINE_CARD
        avatar.line.width = Pt(0.75)
        remove_shadow(avatar)
        # Initial-Letter
        initials = "".join(part[0] for part in name.split() if part[0].isalpha())[:2]
        if not initials:
            initials = "·"
        add_text(slide, initials, ax, card_y + Inches(0.85),
                 Inches(1.6), Inches(0.9),
                 font=FONT_SERIF, size_pt=42, color=MUTED,
                 italic=True, align="center", letter_spacing=-180)
        # Name
        add_text(slide, name, x + Inches(0.4), card_y + Inches(2.3),
                 card_w - Inches(0.8), Inches(0.4),
                 font=FONT_SERIF, size_pt=18, color=TEXT,
                 align="center", letter_spacing=-60)
        # Role
        add_text(slide, role, x + Inches(0.4), card_y + Inches(2.7),
                 card_w - Inches(0.8), Inches(0.4),
                 font=FONT_SANS, size_pt=10, color=TEAL,
                 align="center", letter_spacing=160)
        # Bio
        add_text(slide, bio, x + Inches(0.4), card_y + Inches(3.15),
                 card_w - Inches(0.8), Inches(1.0),
                 font=FONT_SANS, size_pt=10.5, color=TEXT2,
                 align="center", line_spacing=1.45)

    add_footer(slide, 11)


def build_slide_12_ask(prs):
    slide = new_slide(prs)
    add_glow(slide, SLIDE_W // 2, Inches(3.8),
             Inches(11), Inches(6), color=TEAL_GLOW)

    add_eyebrow(slide, "Ask", MARGIN_X, MARGIN_TOP)
    add_serif_heading_dual(
        slide, "Was wir", "suchen.",
        MARGIN_X, Inches(1.2), SLIDE_W - 2 * MARGIN_X,
        size_pt=64, color1=TEXT, color2=TEAL, letter_spacing=-380,
    )

    col_y = Inches(3.5)
    col_h = Inches(3.0)
    col_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.4)) // 2

    # Linke Spalte: Funding + Use of Funds
    left_x = MARGIN_X
    add_card(slide, left_x, col_y, col_w, col_h,
             fill_color=PAPER, border_color=HAIRLINE_CARD, radius_pt=14)
    add_eyebrow(slide, "Funding-Range", left_x + Inches(0.4),
                col_y + Inches(0.35), color=TEAL)
    add_text(slide, "€ X.X M",
             left_x + Inches(0.4), col_y + Inches(0.8),
             col_w - Inches(0.8), Inches(0.8),
             font=FONT_SERIF, size_pt=44, color=TEAL,
             italic=True, letter_spacing=-260)
    add_text(slide, "Pre-Seed / Seed",
             left_x + Inches(0.4), col_y + Inches(1.55),
             col_w - Inches(0.8), Inches(0.3),
             font=FONT_SANS, size_pt=11, color=MUTED, letter_spacing=180)
    # Use of Funds
    use_y = col_y + Inches(2.05)
    add_text(slide, "USE OF FUNDS", left_x + Inches(0.4), use_y,
             col_w - Inches(0.8), Inches(0.3),
             font=FONT_SANS, size_pt=9, color=FAINT,
             bold=True, letter_spacing=220)
    funds = [
        ("Produkt-Tiefe", "AI-Features, Premium-Tier."),
        ("Growth", "Performance Marketing + Content."),
        ("Hiring", "Engineering & GTM."),
    ]
    fy = use_y + Inches(0.35)
    for headline, body in funds:
        add_text(slide, headline, left_x + Inches(0.4), fy,
                 Inches(1.6), Inches(0.3),
                 font=FONT_SANS, size_pt=10, color=TEXT, bold=True)
        add_text(slide, body, left_x + Inches(2.0), fy,
                 col_w - Inches(2.4), Inches(0.3),
                 font=FONT_SANS, size_pt=10, color=TEXT2)
        fy += Inches(0.3)

    # Rechte Spalte: Kontakt
    right_x = MARGIN_X + col_w + Inches(0.4)
    add_card(slide, right_x, col_y, col_w, col_h,
             fill_color=PAPER, border_color=HAIRLINE_CARD, radius_pt=14)
    add_eyebrow(slide, "Kontakt", right_x + Inches(0.4),
                col_y + Inches(0.35), color=TEAL)
    add_text(slide, "Florian Krause",
             right_x + Inches(0.4), col_y + Inches(0.8),
             col_w - Inches(0.8), Inches(0.5),
             font=FONT_SERIF, size_pt=26, color=TEXT,
             letter_spacing=-120)
    add_text(slide, "Founder",
             right_x + Inches(0.4), col_y + Inches(1.25),
             col_w - Inches(0.8), Inches(0.3),
             font=FONT_SANS, size_pt=11, color=MUTED, letter_spacing=180)
    contacts = [
        ("E-MAIL", "florian@tryapplyos.app"),
        ("WEB", "tryapplyos.app"),
    ]
    cy = col_y + Inches(1.75)
    for label, value in contacts:
        add_text(slide, label, right_x + Inches(0.4), cy,
                 Inches(1.2), Inches(0.3),
                 font=FONT_SANS, size_pt=9, color=FAINT,
                 bold=True, letter_spacing=220)
        add_text(slide, value, right_x + Inches(1.6), cy,
                 col_w - Inches(2.0), Inches(0.3),
                 font=FONT_SANS, size_pt=12, color=TEXT)
        cy += Inches(0.4)

    # Logo unten rechts
    if ASSET_LOGO.exists():
        slide.shapes.add_picture(str(ASSET_LOGO),
                                 right_x + col_w - Inches(0.6),
                                 col_y + col_h - Inches(0.65),
                                 height=Inches(0.4))

    # "Danke." zentriert ganz unten
    add_text(slide, "Danke.",
             Inches(0), Inches(6.8), SLIDE_W, Inches(0.5),
             font=FONT_SERIF, size_pt=22, color=TEAL,
             italic=True, align="center", letter_spacing=-80)


# ──────────────────────────────────────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────────────────────────────────────


def main() -> int:
    # Font-Check (informativ)
    print(f"Generating ApplyOS Pitchdeck → {OUTPUT_PATH}")
    print(f"  Repository: {REPO_ROOT}")
    print(f"  Assets:")
    print(f"    logo:      {ASSET_LOGO.exists()} ({ASSET_LOGO})")
    print(f"    result:    {ASSET_RESULT.exists()} ({ASSET_RESULT})")
    print(f"  Fonts: Source Serif 4 + Geist (Fallback: Georgia + Helvetica)")
    print(f"    Hinweis: Falls Fonts nicht installiert sind, nutzt PowerPoint")
    print(f"    automatisch den Fallback. Für beste Optik installieren:")
    print(f"      https://fonts.google.com/specimen/Source+Serif+4")
    print(f"      https://vercel.com/font/geist")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        build_slide_01_cover,
        build_slide_02_problem,
        build_slide_03_solution,
        build_slide_04_market,
        build_slide_05_product_jobsuchende,
        build_slide_06_product_schueler,
        build_slide_07_differentiation,
        build_slide_08_gtm,
        build_slide_09_traction,
        build_slide_10_roadmap,
        build_slide_11_team,
        build_slide_12_ask,
    ]

    for i, build in enumerate(builders, start=1):
        print(f"  → Slide {i:2}/{len(builders)}: {build.__name__}")
        build(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"\n✓ Saved: {OUTPUT_PATH}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Size:   {OUTPUT_PATH.stat().st_size / 1024:.1f} KB")
    return 0


if __name__ == "__main__":
    sys.exit(main())
